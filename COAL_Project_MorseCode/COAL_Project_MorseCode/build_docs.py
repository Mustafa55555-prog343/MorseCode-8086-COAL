"""
build_docs.py
==============
One-shot builder for every document deliverable:

    1. 01_Project_Proposal.docx
    2. 02_Project_Report.docx             (IEEE two-column format)
    3. 03_Presentation.pptx               (13-16 minute deck)
    4. 04_Viva_Preparation_Guide.docx

Just run:  python build_docs.py
"""

from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.enum.section import WD_SECTION, WD_ORIENT
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from docx.shared import Pt, Inches, RGBColor, Cm

from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor as PPTRGBColor


ROOT = Path(__file__).resolve().parent
DOC_DIR = ROOT / "Documents"
DOC_DIR.mkdir(exist_ok=True)


# ==========================================================================
# SHARED HELPERS
# ==========================================================================
def set_cell_background(cell, hex_color: str) -> None:
    """Shade a table cell by injecting <w:shd> directly."""
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), hex_color.lstrip("#"))
    tc_pr.append(shd)


def add_horizontal_line(paragraph) -> None:
    p_pr = paragraph._p.get_or_add_pPr()
    p_bdr = OxmlElement("w:pBdr")
    bottom = OxmlElement("w:bottom")
    bottom.set(qn("w:val"), "single")
    bottom.set(qn("w:sz"), "8")
    bottom.set(qn("w:space"), "1")
    bottom.set(qn("w:color"), "0EA5E9")
    p_bdr.append(bottom)
    p_pr.append(p_bdr)


def styled_run(paragraph, text, *, size=11, bold=False, color=None,
               font="Calibri"):
    run = paragraph.add_run(text)
    run.font.name = font
    rPr = run._element.get_or_add_rPr()
    rFonts = rPr.find(qn("w:rFonts"))
    if rFonts is None:
        rFonts = OxmlElement("w:rFonts")
        rPr.append(rFonts)
    rFonts.set(qn("w:ascii"), font)
    rFonts.set(qn("w:hAnsi"), font)
    run.font.size = Pt(size)
    run.bold = bold
    if color is not None:
        run.font.color.rgb = RGBColor.from_string(color.lstrip("#"))
    return run


def set_default_font(doc: Document, font="Calibri", size=11) -> None:
    style = doc.styles["Normal"]
    style.font.name = font
    style.font.size = Pt(size)
    rPr = style.element.get_or_add_rPr()
    rFonts = rPr.find(qn("w:rFonts"))
    if rFonts is None:
        rFonts = OxmlElement("w:rFonts")
        rPr.append(rFonts)
    rFonts.set(qn("w:ascii"), font)
    rFonts.set(qn("w:hAnsi"), font)


def add_heading(doc, text, *, level=1, color="0F172A",
                size=None, center=False, space_before=6, space_after=3):
    size = size or {1: 20, 2: 15, 3: 13, 4: 12}.get(level, 12)
    p = doc.add_paragraph()
    if center:
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(space_before)
    p.paragraph_format.space_after = Pt(space_after)
    styled_run(p, text, size=size, bold=True, color=color,
               font="Calibri")
    return p


def add_body(doc, text, *, size=11, justify=True,
             space_after=6, color=None, italic=False, bold=False):
    p = doc.add_paragraph()
    if justify:
        p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    p.paragraph_format.space_after = Pt(space_after)
    p.paragraph_format.line_spacing = 1.15
    r = styled_run(p, text, size=size, bold=bold, color=color)
    r.italic = italic
    return p


def add_bullets(doc, items, *, size=11, indent=0.25):
    for it in items:
        p = doc.add_paragraph(style="List Bullet")
        p.paragraph_format.space_after = Pt(3)
        p.paragraph_format.left_indent = Inches(indent)
        styled_run(p, it, size=size)


def ieee_columns(section) -> None:
    """Turn a section into two-column IEEE layout."""
    sect_pr = section._sectPr
    cols = sect_pr.find(qn("w:cols"))
    if cols is None:
        cols = OxmlElement("w:cols")
        sect_pr.append(cols)
    cols.set(qn("w:num"), "2")
    cols.set(qn("w:space"), "480")   # ~0.33"
    cols.set(qn("w:sep"), "0")


def single_column(section) -> None:
    sect_pr = section._sectPr
    cols = sect_pr.find(qn("w:cols"))
    if cols is None:
        cols = OxmlElement("w:cols")
        sect_pr.append(cols)
    cols.set(qn("w:num"), "1")


# ==========================================================================
# 1. PROJECT PROPOSAL
# ==========================================================================
def build_proposal() -> Path:
    doc = Document()
    set_default_font(doc)
    for sec in doc.sections:
        sec.top_margin = Inches(0.8)
        sec.bottom_margin = Inches(0.8)
        sec.left_margin = Inches(0.9)
        sec.right_margin = Inches(0.9)

    # ---- Title block --------------------------------------------------
    t = doc.add_paragraph()
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    styled_run(t, "PROJECT PROPOSAL", size=22, bold=True, color="0F172A")
    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    styled_run(sub, "CS-234  |  Computer Organization & Assembly Language",
               size=12, color="334155")
    sub2 = doc.add_paragraph()
    sub2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    styled_run(sub2, "End-of-Semester Project", size=12, bold=True,
               color="0EA5E9")
    add_horizontal_line(sub2)

    # ---- Meta table ---------------------------------------------------
    tbl = doc.add_table(rows=5, cols=2)
    tbl.autofit = False
    tbl.columns[0].width = Inches(1.8)
    tbl.columns[1].width = Inches(4.5)
    rows = [
        ("Project Title",  "Morse Code Communication and Emergency "
                           "Signaling System"),
        ("Student Name",   "Mustafa Shahid"),
        ("Class",          "BSCS-14B"),
        ("CMS ID",         "500889"),
        ("Faculty",        "Dr. Omar Zeb  |  Sir Ijaz Alam Khan"),
    ]
    for i, (k, v) in enumerate(rows):
        c0, c1 = tbl.rows[i].cells
        c0.text = ""
        c1.text = ""
        styled_run(c0.paragraphs[0], k, size=11, bold=True, color="0F172A")
        styled_run(c1.paragraphs[0], v, size=11)
        set_cell_background(c0, "E2E8F0")

    doc.add_paragraph()

    # ---- Introduction ------------------------------------------------
    add_heading(doc, "1.  Introduction and Description", level=2,
                color="0EA5E9")
    add_body(
        doc,
        "This project implements a Morse Code Communication and Emergency "
        "Signaling System as an end-to-end demonstration of fundamental "
        "Computer Organization and Assembly Language (COAL) concepts.  "
        "Morse code — over 180 years old and still used in aviation, "
        "maritime operations, amateur radio and emergency rescue — remains "
        "a compact, low-bandwidth way to transmit information when modern "
        "channels are unavailable.  The system is delivered in two tightly "
        "coupled layers.  The assembly layer is an 8086 program that "
        "accepts user text through DOS interrupts, converts each character "
        "to its international Morse representation using a fixed-size "
        "lookup table addressed via indexed addressing, and transmits the "
        "result on the PC speaker through BIOS/DOS interrupts.  The "
        "desktop layer is a polished dark-themed Windows application that "
        "mirrors the same logic across fourteen interactive pages: a "
        "welcome/home dashboard with live session statistics, an input-"
        "output translator with sample-phrase chips, history and TXT/WAV "
        "export, a live animated visualiser, an emergency SOS broadcaster, "
        "a hand-operated telegraph key that decodes your press-and-hold "
        "rhythm in real time, an interactive Morse trainer game with "
        "scoring and streaks, a clickable Morse cheat sheet, a CPU "
        "simulator that walks the 8086 through the translation instruction "
        "by instruction with auto-play, speed control and a full status-"
        "flag breakdown, a commented assembly source viewer, a COAL-"
        "concepts reference, an ASCII and Interrupt Vector Table explorer, "
        "a twelve-question interactive quiz, a preferences/settings panel "
        "and an About page.  Together these artifacts cover "
        "instruction set architecture, data definition directives, "
        "arithmetic and logical operations, direct, indirect, indexed "
        "and base-plus-offset addressing modes, software interrupts and "
        "the interrupt vector table, I/O operations, memory-hierarchy "
        "intuition via lookup tables and procedure-based program "
        "structure — giving a compact yet impactful real-world "
        "application that is straightforward to run, explain and defend "
        "during the viva."
    )

    # ---- Roles --------------------------------------------------------
    add_heading(doc, "2.  Roles and Responsibilities",
                level=2, color="0EA5E9")
    add_body(
        doc, "This project is a solo submission.  All roles are undertaken "
        "by the sole member, Mustafa Shahid (BSCS-14B, CMS 500889):",
    )
    add_bullets(doc, [
        "Research and literature survey on Morse code, ITU-R M.1677-1 "
        "timing standards and 8086 BIOS/DOS interrupt services.",
        "Design of the Morse lookup-table data structure and its indexed "
        "addressing scheme.",
        "Implementation of the 8086 assembly program in MASM/emu8086 syntax "
        "with full in-line commentary.",
        "Design and implementation of the Python/CustomTkinter desktop "
        "application (translator, visualiser, SOS broadcaster, assembly "
        "viewer, concept reference).",
        "Technical writing of the proposal, IEEE-format report, "
        "presentation slides and viva-preparation guide.",
        "Testing, polishing, quality-assurance and preparation of the "
        "final demo.",
    ])

    # ---- Deliverables -------------------------------------------------
    add_heading(doc, "3.  Deliverables",
                level=2, color="0EA5E9")
    add_bullets(doc, [
        "Fully-commented 8086 assembly source file (MorseTranslator.asm).",
        "Windows desktop application with a modern dark GUI.",
        "Project report in IEEE two-column format.",
        "Presentation deck prepared for a 13–16-minute talk.",
        "Viva-preparation guide covering all COAL concepts used.",
        "README explaining how to build, run and demonstrate the system.",
    ])

    out = DOC_DIR / "01_Project_Proposal.docx"
    doc.save(out)
    return out


# ==========================================================================
# 2. IEEE PROJECT REPORT
# ==========================================================================
def build_report() -> Path:
    doc = Document()
    set_default_font(doc, size=10)
    for sec in doc.sections:
        sec.top_margin = Inches(0.75)
        sec.bottom_margin = Inches(0.75)
        sec.left_margin = Inches(0.75)
        sec.right_margin = Inches(0.75)

    # ---- TITLE (single-column portion) -------------------------------
    first_section = doc.sections[0]
    single_column(first_section)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    styled_run(
        title,
        "Morse Code Communication and Emergency Signaling System:\n"
        "An 8086 Assembly and Desktop-GUI Implementation",
        size=18, bold=True, color="0F172A",
    )

    author = doc.add_paragraph()
    author.alignment = WD_ALIGN_PARAGRAPH.CENTER
    styled_run(author, "Mustafa Shahid", size=12, bold=True, color="0F172A")
    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    styled_run(
        meta,
        "BSCS-14B  |  CMS ID: 500889\n"
        "CS-234 Computer Organization and Assembly Language\n"
        "Faculty: Dr. Omar Zeb and Sir Ijaz Alam Khan",
        size=10, color="475569",
    )

    # Switch to two-column for the rest of the paper
    doc.add_section(WD_SECTION.CONTINUOUS)
    ieee_columns(doc.sections[-1])

    # ---- Abstract ----------------------------------------------------
    p = doc.add_paragraph()
    styled_run(p, "Abstract — ", size=10, bold=True)
    styled_run(
        p,
        "This paper presents the design and implementation of a Morse code "
        "communication and emergency signaling system developed as a "
        "pedagogical platform for the Computer Organization and Assembly "
        "Language course.  The system comprises a fully commented 8086 "
        "assembly-language program and a modern Python desktop application "
        "with a matching dark-themed graphical user interface.  Core COAL "
        "concepts exercised include indexed, direct, indirect and "
        "base-plus-offset addressing modes; data definition directives; "
        "arithmetic and logical instructions; software interrupts; the "
        "interrupt vector table; procedure calls; and low-level I/O "
        "operations.  The assembly program translates a user-supplied string "
        "into international Morse code by performing O(1) lookups against a "
        "fixed-size record table, while the desktop application demonstrates "
        "the same translation, plays audible dots and dashes at selectable "
        "frequency and duration, visualises transmissions in real time and "
        "includes an emergency SOS broadcaster.  Test runs confirm correct "
        "end-to-end operation, faithful reproduction of ITU-R M.1677-1 "
        "timing and stable behaviour on Windows 11.",
        size=10,
    )

    idx_p = doc.add_paragraph()
    styled_run(idx_p, "Index Terms — ", size=10, bold=True)
    styled_run(
        idx_p,
        "Assembly language, 8086 microprocessor, addressing modes, "
        "interrupts, Morse code, emergency communication, human-computer "
        "interaction, computer organization.",
        size=10,
    )

    # ---- 1. Introduction ---------------------------------------------
    add_heading(doc, "I.  INTRODUCTION", level=2, size=12,
                color="0F172A", space_before=6, space_after=3)
    add_body(
        doc,
        "Samuel Morse’s dot-and-dash alphabet has been in continuous service "
        "since 1844.  Despite the arrival of high-bandwidth digital "
        "protocols, Morse code remains an indispensable fallback channel "
        "for aviation (identification of VORs and NDBs), maritime search "
        "and rescue, military field communications and amateur radio.  Its "
        "enduring relevance rests on three properties: it is slow but "
        "phenomenally robust under poor signal-to-noise ratios; it requires "
        "only an on/off signal source; and it is trivially learnable as an "
        "auditory pattern, making it accessible to operators with "
        "limited hardware.",
        size=10,
    )
    add_body(
        doc,
        "These same properties make Morse code an ideal vehicle for a "
        "concise yet conceptually rich Computer Organization and Assembly "
        "Language (COAL) project.  Converting ASCII characters to Morse "
        "sequences can be expressed as a constant-time lookup against a "
        "fixed-size record table, which in turn motivates a clean use of "
        "indexed addressing.  Signalling the result audibly leverages the "
        "PC-speaker BIOS/DOS interrupt services, which illustrate both the "
        "interrupt vector table and low-level I/O.  Combining translation "
        "and transmission in one compact program therefore exercises the "
        "majority of topics in the CS-234 syllabus.",
        size=10,
    )

    # ---- 2. Literature / Background ---------------------------------
    add_heading(doc, "II.  BACKGROUND AND LITERATURE REVIEW", level=2,
                size=12, color="0F172A")
    add_body(
        doc,
        "The International Telecommunication Union codifies Morse timing in "
        "Recommendation ITU-R M.1677-1 [1].  One \"unit\" is the base "
        "duration; a dot occupies one unit, a dash three, the intra-character "
        "gap is one unit, the inter-character gap three units and the "
        "inter-word gap seven.  Standard character representations are "
        "published in [1] and summarised in most communications textbooks "
        "[2].  On the organisational side, Hennessy and Patterson [3] and "
        "Irvine [4] discuss how small lookup tables map naturally onto "
        "cache-resident memory and how addressing modes allow "
        "constant-time record access.  Intel’s 8086 programmer’s reference "
        "[5] documents the interrupt vector table and the BIOS/DOS services "
        "used in this project.",
        size=10,
    )
    add_body(
        doc,
        "Previous educational implementations of Morse translators tend to "
        "be either (a) pure software demonstrations lacking any "
        "assembly-level exposition, or (b) microcontroller projects that "
        "trade GUI polish for hardware breadth.  This project bridges the "
        "two extremes: a canonical 8086 program showcases the assembly "
        "concepts, while a modern desktop application provides the audible "
        "and visible demonstration required during an in-class viva.",
        size=10,
    )

    # ---- 3. System Architecture -------------------------------------
    add_heading(doc, "III.  SYSTEM ARCHITECTURE", level=2,
                size=12, color="0F172A")
    add_body(
        doc,
        "Fig. 1 shows the overall block diagram.  The system is partitioned "
        "into four layers: (i) user-interface, responsible for input and "
        "output; (ii) translation engine, which performs ASCII-to-Morse "
        "conversion via the lookup table; (iii) signal generator, which "
        "turns dot/dash symbols into audio pulses using the PC speaker or "
        "winsound.Beep; and (iv) the shared data segment that contains the "
        "lookup table and input/output buffers.  All four layers are "
        "present in both the assembly program and the desktop application, "
        "ensuring behavioural parity between the two.",
        size=10,
    )

    # ---- ASCII block diagram -----------------------------------------
    dia = doc.add_paragraph()
    dia.alignment = WD_ALIGN_PARAGRAPH.LEFT
    styled_run(
        dia,
        "+---------------------+       +----------------------+\n"
        "|  User Interface     | <---> |  Translation Engine  |\n"
        "|  (INT 21h / GUI)    |       |  (ASCII -> Morse)    |\n"
        "+---------------------+       +----------------------+\n"
        "           |                              |\n"
        "           v                              v\n"
        "+---------------------+       +----------------------+\n"
        "|  Input/Output Buffer|       |  Morse Lookup Table  |\n"
        "|  (DS:input_buffer)  |       |  (6-byte records)    |\n"
        "+---------------------+       +----------------------+\n"
        "                                         |\n"
        "                                         v\n"
        "                              +----------------------+\n"
        "                              |  Signal Generator    |\n"
        "                              |  (Beep / INT 21h/02) |\n"
        "                              +----------------------+\n",
        size=8, font="Consolas",
    )
    cap = doc.add_paragraph()
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    cap_run = styled_run(cap, "Fig. 1.  System block diagram.",
                         size=9, color="475569")
    cap_run.italic = True

    # ---- 4. Data Structure ------------------------------------------
    add_heading(doc, "IV.  DATA STRUCTURES AND ADDRESSING", level=2,
                size=12, color="0F172A")
    add_body(
        doc,
        "The heart of the system is a single lookup table stored in the data "
        "segment.  Each entry is exactly six bytes: five bytes hold the "
        "Morse representation (dots and dashes, space-padded to uniform "
        "width) and the sixth byte stores the length of the code.  Because "
        "every record is the same size, the address of the entry for "
        "character c is simply:",
        size=10,
    )
    codep = doc.add_paragraph()
    codep.alignment = WD_ALIGN_PARAGRAPH.CENTER
    styled_run(
        codep, "addr = morse_table + (index_of(c) × 6)",
        size=10, bold=True, font="Consolas", color="0EA5E9",
    )
    add_body(
        doc,
        "This reduces Morse lookup to an O(1) operation and maps naturally "
        "onto indexed addressing:  MOV DL, [BX + DI], where BX is the "
        "table base and DI is the computed offset.  The length byte is "
        "reached with the base-plus-index-plus-displacement form MOV CH, "
        "[BX + DI + 5].  Letters are handled by subtracting the ASCII code "
        "of \"A\" to produce an index in 0..25; digits map to slots 26..35 "
        "by subtracting \"0\" and adding 26.  Lowercase characters are "
        "converted to uppercase with a single AND AL, 0DFh — one of the "
        "classical logical-operation tricks that showcases why bit-level "
        "thinking matters in assembly programming.",
        size=10,
    )

    # ---- 5. Control Flow --------------------------------------------
    add_heading(doc, "V.  CONTROL FLOW AND INTERRUPT USAGE", level=2,
                size=12, color="0F172A")
    add_body(
        doc,
        "After initialising DS with the address of the data segment, MAIN "
        "prints a banner via DOS function 09h and solicits input through "
        "buffered read 0Ah.  A translation loop walks SI across the input, "
        "applies uppercase normalisation, resolves the table entry as "
        "described above, emits the Morse symbols through DOS function 02h "
        "(single-character display) and inserts a one-space gap between "
        "letters.  Unsupported characters are silently skipped — a "
        "conscious robustness choice that prevents malformed input from "
        "halting the demo.  At the end of the loop, PLAY_SOS is called "
        "unconditionally to provide an audible demonstration of the "
        "system’s emergency capability.  Finally the program terminates "
        "cleanly via DOS function 4Ch.",
        size=10,
    )
    add_body(
        doc,
        "All I/O is routed through software interrupts, which demonstrates "
        "the role of the interrupt vector table (IVT).  When an INT n "
        "instruction executes, the CPU consults slot n of the IVT at "
        "physical address 0000:n×4 to locate the handler to run.  This "
        "project exercises INT 21h (DOS services: 09h, 0Ah, 02h, 4Ch) as "
        "well as generating ASCII BEL (07h) to invoke the speaker.",
        size=10,
    )

    # ---- 6. Desktop Application -------------------------------------
    add_heading(doc, "VI.  DESKTOP APPLICATION", level=2,
                size=12, color="0F172A")
    add_body(
        doc,
        "To produce an effective demonstration, a companion Windows "
        "desktop application was implemented in Python using "
        "CustomTkinter.  The application is partitioned into fourteen "
        "navigable pages backed by a scrollable sidebar: Home, "
        "Translator, Live Visualiser, Emergency SOS, Telegraph Key, "
        "Morse Trainer, Morse Cheat Sheet, CPU Simulator, Assembly "
        "Code, COAL Concepts, ASCII and IVT Reference, COAL Quiz, "
        "Settings and About.  The Home page acts as a dashboard with "
        "one-click buttons, a live session-statistics strip (number of "
        "translations, audio transmissions, simulator steps and best "
        "trainer score) and a grid of accent-striped cards summarising "
        "each module.  The Translator supports live bidirectional "
        "encoding, a row of sample-phrase chips for one-click demos, "
        "clipboard copy, opening .txt files for batch translation, a "
        "scrollable session-history panel and export of the current "
        "Morse both as a plain-text file and as a .wav audio recording "
        "generated directly from sine-wave samples at the currently "
        "selected frequency and unit duration.  Tone frequency (300–1500 "
        "Hz) and unit duration (50–250 ms) are controlled via sliders.  "
        "The Visualiser flashes a dot or dash on a canvas synchronously "
        "with each audible beep, and the Emergency SOS page triggers a "
        "repeating distress broadcast with nine illuminated indicator "
        "lamps that ignite in sequence with the transmitted elements.",
        size=10,
    )
    add_body(
        doc,
        "Several additional pages exist to deepen the experience.  The "
        "Telegraph Key page lets the user tap out Morse by hand: "
        "pressing the large button (or the SPACE key) measures the "
        "hold duration against the current unit setting and classifies "
        "each tap as a dot or a dash, commits a letter after a 3-unit "
        "gap and a word after a 7-unit gap, decoding the operator's "
        "rhythm into plain text in real time.  The Morse Trainer "
        "gamifies learning: a random letter, word or phrase is shown "
        "as Morse symbols, the user types the plain text, and the "
        "application awards points based on accuracy and speed, with a "
        "streak multiplier and a persistent best-score.  The Morse "
        "Cheat Sheet displays every supported character as a clickable "
        "card which, when clicked, audibly transmits its Morse code.  "
        "The CPU Simulator is described separately in Section VII.  "
        "The COAL Concepts page exposes fourteen reference cards, each "
        "showing a topic, an assembly example and a plain-language "
        "explanation.  The ASCII and IVT Reference contains a "
        "searchable character table with Morse mappings alongside the "
        "Interrupt Vector Table entries invoked by the program.  The "
        "Quiz page offers twelve multiple-choice questions with instant "
        "feedback.  Finally the Settings page persists user preferences "
        "(theme variant, default frequency and unit, audio on/off, "
        "first-run flag) to disk.  The application starts maximised "
        "on Windows, resizes cleanly to any laptop screen, shows "
        "contextual tooltips on every control and includes a first-run "
        "guided welcome overlay plus an F1 keyboard-shortcut dialog.",
        size=10,
    )

    # ---- 7. CPU Simulator --------------------------------------------
    add_heading(doc, "VII.  8086 CPU SIMULATOR", level=2,
                size=12, color="0F172A")
    add_body(
        doc,
        "The CPU Simulator page turns the desktop application into a "
        "pedagogical microscope.  The user types any letter or digit and "
        "clicks Load; the simulator then constructs a trace of the exact "
        "instructions the assembly program would issue to translate that "
        "character.  Pressing Step advances the simulator one instruction "
        "at a time.  A two-column register panel displays AX, BX, CX, "
        "DX, SI, DI, SP, BP, IP and FLAGS as four-digit hexadecimal "
        "values; registers modified by the current instruction are "
        "highlighted in cyan (and IP in amber), mimicking the visual "
        "feedback of professional debuggers.  A memory view lists the "
        "full 216-byte Morse lookup table with an arrow marker pointing "
        "at the record being accessed, visually reinforcing the "
        "index = character × 6 relationship.  A large instruction "
        "display shows the current MOV / SUB / MUL / INT line in a "
        "monospace font, with a natural-language explanation below it.",
        size=10,
    )
    add_body(
        doc,
        "The Simulator was designed as a viva-friendly aid: examiners "
        "can ask the student to translate any character and watch "
        "register-by-register how the computation proceeds — from the "
        "SUB AL,'A' that converts ASCII to a zero-based index, through "
        "the MUL DL that multiplies by the record size, to the "
        "MOV DL, [BX + DI] that performs the actual Morse lookup using "
        "indexed addressing.  In addition to single-step execution, a "
        "Play / Pause button drives the simulator automatically at a "
        "speed chosen with a slider (150–1500 ms between instructions), "
        "and a dedicated FLAGS panel breaks the status register into "
        "its six individual bits — Zero, Carry, Sign, Overflow, Parity "
        "and Auxiliary — each illuminated when set.  A Run-All button "
        "fast-forwards the whole trace, while Reset returns the "
        "simulator to a clean state.",
        size=10,
    )

    # ---- 8. Results --------------------------------------------------
    add_heading(doc, "VIII.  RESULTS AND TESTING", level=2,
                size=12, color="0F172A")
    add_body(
        doc,
        "A battery of representative inputs was run against both the "
        "assembly program and the desktop application, and the outputs "
        "agree exactly.  Table I summarises a subset.",
        size=10,
    )

    tbl = doc.add_table(rows=5, cols=2)
    tbl.style = "Light Grid Accent 1"
    headers = [("Input", "Expected Morse Output"),
               ("SOS",            "... --- ..."),
               ("HELLO WORLD",    ".... . .-.. .-.. --- / .-- --- .-. .-.. -.."),
               ("BSCS 14B",       "-... ... -.-. ... / .---- ....- -..."),
               ("500889",         "..... ----- ----- ---.. ---.. ----.")]
    for i, (a, b) in enumerate(headers):
        c0, c1 = tbl.rows[i].cells
        c0.text = ""
        c1.text = ""
        r0 = styled_run(c0.paragraphs[0], a, size=9,
                        bold=(i == 0), font="Consolas")
        r1 = styled_run(c1.paragraphs[0], b, size=9,
                        bold=(i == 0), font="Consolas")
        if i == 0:
            set_cell_background(c0, "0EA5E9")
            set_cell_background(c1, "0EA5E9")
            r0.font.color.rgb = RGBColor.from_string("FFFFFF")
            r1.font.color.rgb = RGBColor.from_string("FFFFFF")

    cap = doc.add_paragraph()
    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = styled_run(cap, "TABLE I.  Selected translation results.",
                   size=9, color="475569")
    r.italic = True

    add_body(
        doc,
        "The system was further evaluated by broadcasting SOS repeatedly "
        "and by changing frequency and unit duration at run time; all "
        "configurations produced audibly correct transmissions.  No "
        "run-time errors, hangs or visual artefacts were observed.",
        size=10,
    )

    # ---- 9. Opportunities & Challenges ------------------------------
    add_heading(doc, "IX.  OPPORTUNITIES AND CHALLENGES", level=2,
                size=12, color="0F172A")
    add_body(
        doc,
        "Opportunities.  A low-cost, battery-friendly Morse terminal is "
        "immediately useful in disaster-response kits, marine lifeboats, "
        "remote pilot training and amateur radio.  Because the encoding is "
        "so compact, the same firmware can be deployed to an 8-bit "
        "microcontroller to build a physical SOS beacon powered by a "
        "coin cell.  A second opportunity lies in accessibility: Morse "
        "output is a proven communication modality for users with "
        "profound speech or motor impairments.",
        size=10,
    )
    add_body(
        doc,
        "Challenges.  The assembly implementation must balance portability "
        "(MASM, TASM and emu8086 differ in minor dialectal issues) with "
        "clarity.  Audible signalling via INT 21h/02h + BEL is "
        "coarse-grained; producing cleanly tuned tones requires either "
        "direct port programming (ports 0x61 and 0x43 on legacy hardware) "
        "or, as adopted here, a higher-level API such as winsound.Beep.  "
        "Finally, robust handling of unsupported characters and "
        "internationalisation beyond A-Z/0-9 is deliberately kept simple "
        "to preserve the pedagogical clarity of the code.",
        size=10,
    )

    # ---- 10. Ethical Perspective -------------------------------------
    add_heading(doc, "X.  ETHICAL PERSPECTIVE", level=2,
                size=12, color="0F172A")
    add_body(
        doc,
        "Any communication tool carries dual-use potential.  A Morse "
        "transmitter can save lives by broadcasting a distress call from a "
        "crippled vessel, or it can be misused to exchange covert "
        "information in environments where discretion is unwarranted.  "
        "Because the system is open, inspectable and uses a globally "
        "agreed standard, its misuse risk is low — but ethical deployment "
        "requires operators to respect frequency allocations, to identify "
        "themselves on air and never to transmit false distress calls, "
        "which are criminal offences under most national communications "
        "acts.  The project also respects intellectual-property norms: "
        "all code is original, and cited materials are listed in the "
        "references.",
        size=10,
    )

    # ---- 11. Conclusion & Future Work -------------------------------
    add_heading(doc, "XI.  CONCLUSION AND FUTURE WORK",
                level=2, size=12, color="0F172A")
    add_body(
        doc,
        "A Morse code communication and emergency signaling system has "
        "been designed, implemented and tested at three levels of "
        "abstraction.  The 8086 assembly program exercises every major "
        "COAL topic of the course; a step-through CPU simulator exposes "
        "that behaviour register-by-register inside the GUI; and the "
        "companion desktop application provides an auditorium-ready demo "
        "with audible output, an animated visualiser, an emergency SOS "
        "broadcaster, a clickable Morse cheat sheet, an ASCII and IVT "
        "reference and a self-assessment quiz.  Future work could extend "
        "the hardware track to an Arduino or STM32 beacon with an LED "
        "indicator, add bidirectional radio with a USB sound card, and "
        "integrate a trainable audio decoder that could accept Morse "
        "input by microphone — turning the project into a complete "
        "amateur-radio assistant.",
        size=10,
    )

    # ---- References --------------------------------------------------
    add_heading(doc, "REFERENCES", level=2, size=12, color="0F172A")
    refs = [
        "[1]  International Telecommunication Union, Recommendation ITU-R "
        "M.1677-1, \"International Morse code,\" 2009.",
        "[2]  J. G. Proakis and M. Salehi, Communication Systems "
        "Engineering, 2nd ed.  Upper Saddle River, NJ: Prentice Hall, 2002.",
        "[3]  J. L. Hennessy and D. A. Patterson, Computer Organization and "
        "Design: The Hardware/Software Interface, 6th ed.  Cambridge, MA: "
        "Morgan Kaufmann, 2020.",
        "[4]  K. R. Irvine, Assembly Language for x86 Processors, 8th ed.  "
        "Hoboken, NJ: Pearson, 2019.",
        "[5]  Intel Corporation, 8086 Family User’s Manual.  Santa Clara, "
        "CA: Intel Corp., 1979.",
        "[6]  M. Abrash, Zen of Assembly Language: Volume I - Knowledge.  "
        "Glenview, IL: Scott, Foresman, 1990.",
        "[7]  American Radio Relay League, The ARRL Handbook for Radio "
        "Communications, 99th ed., Newington, CT, 2022.",
    ]
    for r in refs:
        p = doc.add_paragraph()
        p.paragraph_format.left_indent = Inches(0.18)
        p.paragraph_format.first_line_indent = Inches(-0.18)
        p.paragraph_format.space_after = Pt(2)
        styled_run(p, r, size=9)

    out = DOC_DIR / "02_Project_Report.docx"
    doc.save(out)
    return out


# ==========================================================================
# 3. POWERPOINT PRESENTATION
# ==========================================================================
NAVY       = PPTRGBColor(0x0F, 0x17, 0x2A)
PANEL      = PPTRGBColor(0x1E, 0x29, 0x3B)
CYAN       = PPTRGBColor(0x38, 0xBD, 0xF8)
CYAN_DIM   = PPTRGBColor(0x0E, 0xA5, 0xE9)
TEXT_LIGHT = PPTRGBColor(0xF1, 0xF5, 0xF9)
TEXT_DIM   = PPTRGBColor(0x94, 0xA3, 0xB8)
ACCENT_GREEN = PPTRGBColor(0x22, 0xC5, 0x5E)
ACCENT_RED   = PPTRGBColor(0xEF, 0x44, 0x44)


def _slide_bg(slide, color=NAVY) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, x, y, w, h, color, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if not line:
        shape.line.fill.background()
    return shape


def _add_text(slide, text, x, y, w, h, *, size=18, bold=False,
              color=TEXT_LIGHT, align=PP_ALIGN.LEFT, font="Calibri",
              anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top  = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = PPt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def _add_bullets(slide, items, x, y, w, h, *, size=16, color=TEXT_LIGHT,
                 bullet_color=CYAN, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = PPt(6)
        r1 = p.add_run()
        r1.text = "▪  "
        r1.font.name = font
        r1.font.size = PPt(size)
        r1.font.bold = True
        r1.font.color.rgb = bullet_color
        r2 = p.add_run()
        r2.text = text
        r2.font.name = font
        r2.font.size = PPt(size)
        r2.font.color.rgb = color
    return tb


def _title_strip(slide, title, subtitle=None):
    _add_rect(slide, PInches(0), PInches(0),
              PInches(13.333), PInches(0.9), PANEL)
    _add_rect(slide, PInches(0), PInches(0.88),
              PInches(13.333), PInches(0.04), CYAN)
    _add_text(slide, title, PInches(0.55), PInches(0.18),
              PInches(10.8), PInches(0.55),
              size=24, bold=True, color=TEXT_LIGHT,
              anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        _add_text(slide, subtitle, PInches(0.55), PInches(0.52),
                  PInches(10.8), PInches(0.4),
                  size=12, color=TEXT_DIM, anchor=MSO_ANCHOR.MIDDLE)
    _add_text(slide, "BSCS-14B  |  500889",
              PInches(11.3), PInches(0.30),
              PInches(2.0), PInches(0.4),
              size=11, bold=True, color=CYAN, align=PP_ALIGN.RIGHT)


def _footer(slide, number, total):
    _add_text(slide, f"Mustafa Shahid   ·   CS-234 COAL Project",
              PInches(0.55), PInches(7.10),
              PInches(8.0), PInches(0.3),
              size=10, color=TEXT_DIM)
    _add_text(slide, f"{number} / {total}",
              PInches(12.3), PInches(7.10),
              PInches(1.0), PInches(0.3),
              size=10, color=TEXT_DIM, align=PP_ALIGN.RIGHT)


def _card(slide, x, y, w, h, title, body, *, accent=CYAN):
    _add_rect(slide, x, y, w, h, PANEL)
    _add_rect(slide, x, y, PInches(0.1), h, accent)
    _add_text(slide, title, x + PInches(0.25), y + PInches(0.15),
              w - PInches(0.35), PInches(0.4),
              size=16, bold=True, color=TEXT_LIGHT)
    _add_text(slide, body, x + PInches(0.25), y + PInches(0.55),
              w - PInches(0.35), h - PInches(0.65),
              size=12, color=TEXT_DIM)


def build_pptx() -> Path:
    prs = Presentation()
    prs.slide_width  = PInches(13.333)
    prs.slide_height = PInches(7.5)

    blank = prs.slide_layouts[6]
    total_slides = 18

    # ----- Slide 1: Title ---------------------------------------------
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _add_rect(s, PInches(0), PInches(3.1), PInches(13.333),
              PInches(0.06), CYAN)

    _add_text(s, "MORSE CODE",
              PInches(0.8), PInches(1.6), PInches(12), PInches(1.2),
              size=64, bold=True, color=TEXT_LIGHT,
              align=PP_ALIGN.CENTER)
    _add_text(s, "Communication & Emergency Signaling System",
              PInches(0.8), PInches(2.35), PInches(12), PInches(0.8),
              size=26, color=CYAN, align=PP_ALIGN.CENTER)
    _add_text(s, "An 8086 Assembly + Desktop-GUI Implementation",
              PInches(0.8), PInches(3.35), PInches(12), PInches(0.55),
              size=18, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    _add_text(s, "Mustafa Shahid",
              PInches(0.8), PInches(5.2), PInches(12), PInches(0.5),
              size=24, bold=True, color=TEXT_LIGHT,
              align=PP_ALIGN.CENTER)
    _add_text(s, "BSCS-14B   |   CMS ID: 500889",
              PInches(0.8), PInches(5.7), PInches(12), PInches(0.4),
              size=16, color=TEXT_DIM, align=PP_ALIGN.CENTER)
    _add_text(s, "CS-234  Computer Organization & Assembly Language",
              PInches(0.8), PInches(6.1), PInches(12), PInches(0.4),
              size=14, color=TEXT_DIM, align=PP_ALIGN.CENTER)
    _add_text(s, "Faculty:  Dr. Omar Zeb  ·  Sir Ijaz Alam Khan",
              PInches(0.8), PInches(6.45), PInches(12), PInches(0.4),
              size=13, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    # ----- Slide 2: Outline -------------------------------------------
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _title_strip(s, "Presentation Outline",
                 "A roadmap for the next 13–16 minutes")
    items = [
        "Why Morse? — motivation and real-world impact",
        "Problem statement and project objectives",
        "Literature review and background",
        "System architecture and block diagram",
        "Data structures and addressing modes",
        "Assembly walk-through (8086)",
        "Software interrupts and the IVT",
        "Desktop application — fourteen interactive pages",
        "Telegraph Key, Trainer and Cheat Sheet",
        "CPU simulator — registers, memory, flags, auto-play",
        "Live demonstration",
        "COAL concepts demonstrated",
        "Opportunities, challenges and ethics",
        "Conclusion and Q&A",
    ]
    _add_bullets(s, items, PInches(1.0), PInches(1.4),
                 PInches(11.0), PInches(5.5), size=18)
    _footer(s, 2, total_slides)

    # ----- Slide 3: Motivation ----------------------------------------
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _title_strip(s, "Why Morse Code?",
                 "A 180-year-old idea that still saves lives")
    _card(s, PInches(0.7), PInches(1.4), PInches(6.0), PInches(2.5),
          "Real-world impact",
          "Aviation beacons (VOR/NDB IDs), maritime distress, "
          "military field-comms, amateur radio, disaster response "
          "and accessibility for speech-impaired users.")
    _card(s, PInches(6.9), PInches(1.4), PInches(5.7), PInches(2.5),
          "Why it still wins",
          "Needs only on/off signalling, survives extreme noise, "
          "works with milliwatts of power and is legible by the "
          "human ear alone.")
    _card(s, PInches(0.7), PInches(4.1), PInches(6.0), PInches(2.5),
          "Why it is perfect for COAL",
          "Fixed-size lookup table ⇒ indexed addressing; I/O via "
          "INT 21h; tight loops, arithmetic and logical tricks "
          "such as AND AL, 0DFh for case conversion.",
          accent=ACCENT_GREEN)
    _card(s, PInches(6.9), PInches(4.1), PInches(5.7), PInches(2.5),
          "What this project delivers",
          "Full 8086 program, modern Windows desktop GUI with "
          "audible beeps, live visualiser, emergency SOS and an "
          "in-app COAL concepts reference.",
          accent=ACCENT_RED)
    _footer(s, 3, total_slides)

    # ----- Slide 4: Problem Statement ---------------------------------
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _title_strip(s, "Problem Statement & Objectives")
    _add_text(s, "Problem",
              PInches(0.7), PInches(1.3), PInches(12), PInches(0.4),
              size=20, bold=True, color=CYAN)
    _add_text(s,
              "Build a compact COAL artefact that translates plain text "
              "into international Morse code, transmits it audibly and "
              "transparently demonstrates the major assembly-language "
              "concepts of the CS-234 syllabus.",
              PInches(0.7), PInches(1.75), PInches(12), PInches(1.2),
              size=16, color=TEXT_LIGHT)

    _add_text(s, "Objectives",
              PInches(0.7), PInches(3.1), PInches(12), PInches(0.4),
              size=20, bold=True, color=CYAN)
    _add_bullets(s, [
        "Implement ASCII → Morse in 8086 assembly using a lookup table.",
        "Demonstrate direct, indirect, indexed and base+offset addressing.",
        "Use BIOS/DOS software interrupts for all I/O.",
        "Build a polished Windows GUI that mirrors the assembly behaviour.",
        "Include an emergency SOS broadcaster and a live visualiser.",
        "Deliver IEEE report, slides, viva guide and a runnable demo.",
    ], PInches(0.7), PInches(3.6), PInches(12), PInches(3.5), size=16)
    _footer(s, 4, total_slides)

    # ----- Slide 5: Literature Review ---------------------------------
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _title_strip(s, "Literature & Background")
    _card(s, PInches(0.7), PInches(1.3), PInches(5.9), PInches(2.7),
          "ITU-R M.1677-1",
          "The international standard for Morse timing: unit, dot, "
          "dash and inter-element gaps in ratios 1 : 1 : 3 : 1 : 3 : 7.")
    _card(s, PInches(6.8), PInches(1.3), PInches(5.9), PInches(2.7),
          "Hennessy & Patterson",
          "Lookup tables are a canonical way to exploit locality of "
          "reference and map onto cache-friendly memory regions.")
    _card(s, PInches(0.7), PInches(4.1), PInches(5.9), PInches(2.7),
          "Irvine (x86 Assembly)",
          "Reference text on addressing modes, interrupt usage and "
          "MASM programming style, followed closely in our code.",
          accent=ACCENT_GREEN)
    _card(s, PInches(6.8), PInches(4.1), PInches(5.9), PInches(2.7),
          "Intel 8086 Manual",
          "Authoritative description of the ISA, the interrupt "
          "vector table and the DOS/BIOS service numbers used here.",
          accent=ACCENT_RED)
    _footer(s, 5, total_slides)

    # ----- Slide 6: Architecture --------------------------------------
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _title_strip(s, "System Architecture", "Four layers, shared memory")
    ascii_art = ("+---------------------+       +----------------------+\n"
                 "|  User Interface     | <---> |  Translation Engine  |\n"
                 "|  (INT 21h / GUI)    |       |  (ASCII -> Morse)    |\n"
                 "+---------------------+       +----------------------+\n"
                 "           |                              |\n"
                 "           v                              v\n"
                 "+---------------------+       +----------------------+\n"
                 "|  Input/Output Buffer|       |  Morse Lookup Table  |\n"
                 "|  (DS:input_buffer)  |       |  (6-byte records)    |\n"
                 "+---------------------+       +----------------------+\n"
                 "                                         |\n"
                 "                                         v\n"
                 "                              +----------------------+\n"
                 "                              |  Signal Generator    |\n"
                 "                              |  (Beep / INT 21h/02) |\n"
                 "                              +----------------------+")
    tb = _add_text(s, ascii_art,
                   PInches(1.4), PInches(1.35),
                   PInches(10.6), PInches(5.4),
                   size=14, color=TEXT_LIGHT, font="Consolas")
    _footer(s, 6, total_slides)

    # ----- Slide 7: Data Structure -----------------------------------
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _title_strip(s, "Lookup Table Layout",
                 "Six bytes per record  →  O(1) indexed access")
    _add_text(s, "Record:   [ . . . . . ] [len]",
              PInches(0.7), PInches(1.35), PInches(12), PInches(0.5),
              size=22, bold=True, color=CYAN, font="Consolas")
    _add_text(s,
              "byte 0-4 : Morse symbols (dots & dashes, space-padded)\n"
              "byte 5   : code length (1 – 5)",
              PInches(0.7), PInches(1.9), PInches(12), PInches(1.2),
              size=16, color=TEXT_DIM, font="Consolas")
    _add_text(s, "addr = morse_table + (index × 6)",
              PInches(0.7), PInches(3.3), PInches(12), PInches(0.6),
              size=24, bold=True, color=ACCENT_GREEN, font="Consolas")
    _add_bullets(s, [
        "Letters  'A' .. 'Z'  →  index 0 .. 25",
        "Digits   '0' .. '9'  →  index 26 .. 35",
        "Case conversion:  AND AL, 0DFh  (clears bit-5)",
        "Letter address fetch:  MOV DL, [BX + DI]",
        "Length fetch:          MOV CH, [BX + DI + 5]",
    ], PInches(0.7), PInches(4.1), PInches(12), PInches(2.7), size=17)
    _footer(s, 7, total_slides)

    # ----- Slide 8: Addressing Modes ---------------------------------
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _title_strip(s, "Addressing Modes in Action")
    cards = [
        ("Direct",   "MOV CL, [input_buffer+1]",
                     "Named absolute location — reads length field."),
        ("Indirect", "MOV AL, [SI]",
                     "SI points to the current input character."),
        ("Indexed",  "MOV AL, [SI + 2]",
                     "Base pointer SI + constant displacement."),
        ("Base+Index","MOV DL, [BX + DI]",
                     "BX = table base, DI = index × 6."),
        ("Base+Idx+Disp","MOV CH, [BX + DI + 5]",
                     "Jumps to the length byte of the record."),
        ("Logical",  "AND AL, 0DFh",
                     "Bitmask-based lowercase → uppercase."),
    ]
    positions = [
        (PInches(0.6),  PInches(1.3), PInches(6.05), PInches(1.65)),
        (PInches(6.75), PInches(1.3), PInches(6.05), PInches(1.65)),
        (PInches(0.6),  PInches(3.05),PInches(6.05), PInches(1.65)),
        (PInches(6.75), PInches(3.05),PInches(6.05), PInches(1.65)),
        (PInches(0.6),  PInches(4.8), PInches(6.05), PInches(1.65)),
        (PInches(6.75), PInches(4.8), PInches(6.05), PInches(1.65)),
    ]
    accents = [CYAN, CYAN, CYAN_DIM, CYAN_DIM, ACCENT_GREEN, ACCENT_RED]
    for (mode, code, expl), (x, y, w, h), acc in zip(cards, positions, accents):
        _add_rect(s, x, y, w, h, PANEL)
        _add_rect(s, x, y, PInches(0.1), h, acc)
        _add_text(s, mode, x + PInches(0.25), y + PInches(0.1),
                  w - PInches(0.35), PInches(0.4),
                  size=15, bold=True, color=TEXT_LIGHT)
        _add_text(s, code, x + PInches(0.25), y + PInches(0.5),
                  w - PInches(0.35), PInches(0.4),
                  size=13, bold=True, color=CYAN, font="Consolas")
        _add_text(s, expl, x + PInches(0.25), y + PInches(0.9),
                  w - PInches(0.35), PInches(0.7),
                  size=12, color=TEXT_DIM)
    _footer(s, 8, total_slides)

    # ----- Slide 9: Assembly Walkthrough ------------------------------
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _title_strip(s, "Assembly Walk-through",
                 "The flow of MAIN at a glance")
    _add_bullets(s, [
        "Initialize DS with the data segment address.",
        "Print banner and prompt via DOS function 09h.",
        "Read user text with buffered input — DOS function 0Ah.",
        "For each character: normalise case with AND AL, 0DFh.",
        "Compute index — digits: SUB + ADD 26; letters: SUB 'A'.",
        "Multiply by RECORD_SIZE (6) to get the table offset.",
        "Loop CH times, emitting each Morse symbol with function 02h.",
        "Invoke PLAY_SOS — audible emergency signal via BEL (07h).",
        "Terminate cleanly with DOS function 4Ch.",
    ], PInches(0.7), PInches(1.4), PInches(12), PInches(5.5), size=17)
    _footer(s, 9, total_slides)

    # ----- Slide 10: Interrupts ---------------------------------------
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _title_strip(s, "Software Interrupts & the IVT")
    _add_text(s,
              "INT n  →  consult IVT slot n at 0000:n×4  →  run OS handler",
              PInches(0.7), PInches(1.35), PInches(12), PInches(0.6),
              size=18, bold=True, color=CYAN, font="Consolas")

    rows = [
        ("INT 21h / 09h", "Display $-terminated string"),
        ("INT 21h / 0Ah", "Buffered keyboard input"),
        ("INT 21h / 02h", "Display single character"),
        ("INT 21h / 4Ch", "Terminate with exit code"),
        ("ASCII 07h",     "BEL — PC-speaker beep"),
    ]
    y0 = PInches(2.1)
    for i, (left, right) in enumerate(rows):
        y = y0 + PInches(0.7 * i)
        _add_rect(s, PInches(0.7), y, PInches(11.8), PInches(0.62), PANEL)
        _add_rect(s, PInches(0.7), y, PInches(0.08), PInches(0.62), CYAN)
        _add_text(s, left, PInches(0.9), y + PInches(0.1),
                  PInches(3.2), PInches(0.45),
                  size=15, bold=True, color=CYAN, font="Consolas",
                  anchor=MSO_ANCHOR.MIDDLE)
        _add_text(s, right, PInches(4.2), y + PInches(0.1),
                  PInches(8.2), PInches(0.45),
                  size=15, color=TEXT_LIGHT,
                  anchor=MSO_ANCHOR.MIDDLE)
    _footer(s, 10, total_slides)

    # ----- Slide 11: Desktop App --------------------------------------
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _title_strip(s, "Desktop Application",
                 "Python + CustomTkinter, fourteen interactive pages")
    _add_bullets(s, [
        "Home dashboard — live session statistics + one-click cards.",
        "Translator:  live ASCII ↔ Morse, TXT/WAV export, sample chips.",
        "Live Visualiser:  big dot/dash flashes synced to audio.",
        "Emergency SOS:  nine indicator lamps + broadcast button.",
        "Telegraph Key:  press & hold to tap Morse by hand.",
        "Morse Trainer:  decode challenges with score + streak.",
        "Cheat Sheet:  click any card to hear that symbol.",
        "CPU Simulator:  step, auto-play, speed slider, flag breakdown.",
        "Assembly Viewer, COAL Concepts, ASCII/IVT, 12-q Quiz, Settings.",
        "Starts maximised; responsive from 960×620 to 4K.",
    ], PInches(0.7), PInches(1.4), PInches(12), PInches(5.5), size=15)
    _footer(s, 11, total_slides)

    # ----- Slide 12: CPU Simulator ------------------------------------
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _title_strip(s, "CPU Simulator",
                 "Watch the 8086 execute — register by register")
    _card(s, PInches(0.7), PInches(1.3), PInches(6.0), PInches(2.5),
          "What it does",
          "Type any character, click Load. The simulator builds the "
          "exact instruction trace the assembly code would execute, "
          "ready to step through.")
    _card(s, PInches(6.9), PInches(1.3), PInches(5.7), PInches(2.5),
          "What lights up",
          "AX, BX, CX, DX, SI, DI, SP, BP, IP and FLAGS are shown as "
          "hex values.  Cyan = modified this step.  Amber = IP advance.",
          accent=ACCENT_GREEN)
    _card(s, PInches(0.7), PInches(4.0), PInches(6.0), PInches(2.7),
          "Memory view",
          "The full 216-byte Morse lookup table is rendered with an "
          "arrow marker that points at  base + (index × 6)  during "
          "each indexed fetch.",
          accent=ACCENT_RED)
    _card(s, PInches(6.9), PInches(4.0), PInches(5.7), PInches(2.7),
          "Pedagogical win",
          "Students — and examiners — can inspect every step of the "
          "translation in natural language, directly inside the demo, "
          "without touching emu8086.")
    _footer(s, 12, total_slides)

    # ----- Slide 13: Live Demo ----------------------------------------
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _title_strip(s, "Live Demonstration")
    _add_text(s,
              "Demo script (≈4 minutes):",
              PInches(0.7), PInches(1.35), PInches(12), PInches(0.5),
              size=20, bold=True, color=CYAN)
    _add_bullets(s, [
        "1.  Double-click Run_Application.bat  →  window opens maximised.",
        "2.  Home  →  show statistics strip + click a Translator chip.",
        "3.  Translator  →  type HELLO WORLD  →  ▶ Play Audio.",
        "4.  Save .txt and .wav  →  show the files on disk.",
        "5.  Cheat Sheet  →  click a few cards  →  hear each letter.",
        "6.  Telegraph Key  →  hold SPACE to tap SOS by hand.",
        "7.  Morse Trainer  →  answer one challenge  →  show score.",
        "8.  CPU Simulator  →  type 'S', Load, Play at medium speed.",
        "9.  Live Visualiser  →  transmit SOS HELP  →  watch flashes.",
        "10. Emergency SOS  →  BROADCAST SOS  →  lamps sequence.",
        "11. ASCII/IVT + COAL Concepts + Quiz  →  walk the cards.",
        "12. Settings  →  switch theme  →  show preferences persist.",
    ], PInches(0.7), PInches(1.9), PInches(12), PInches(5.2), size=15)
    _footer(s, 13, total_slides)

    # ----- Slide 14: COAL Concepts ------------------------------------
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _title_strip(s, "COAL Concepts Demonstrated")

    concepts = [
        "Instruction Set Architecture (8086)",
        "Data Definition Directives (DB, EQU)",
        "Direct Addressing Mode",
        "Indirect Addressing Mode",
        "Indexed Addressing Mode",
        "Base + Index + Displacement",
        "Arithmetic & Logical Operations",
        "Software Interrupts (INT 21h)",
        "Interrupt Vector Table",
        "Procedures, CALL / RET, the Stack",
        "I/O Operations (PC Speaker)",
        "Memory Hierarchy via Lookup Tables",
    ]
    x_left, x_right = PInches(0.7), PInches(6.9)
    y0 = PInches(1.35)
    for i, c in enumerate(concepts):
        col = i % 2
        row = i // 2
        x = x_left if col == 0 else x_right
        y = y0 + PInches(row * 0.85)
        _add_rect(s, x, y, PInches(5.9), PInches(0.75), PANEL)
        _add_rect(s, x, y, PInches(0.08), PInches(0.75), CYAN)
        _add_text(s, c, x + PInches(0.22), y + PInches(0.12),
                  PInches(5.6), PInches(0.55),
                  size=15, bold=True, color=TEXT_LIGHT,
                  anchor=MSO_ANCHOR.MIDDLE)
    _footer(s, 14, total_slides)

    # ----- Slide 15: Feature Showcase (new pages) ---------------------
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _title_strip(s, "Feature Showcase",
                 "What the desktop application delivers")
    _card(s, PInches(0.6), PInches(1.3), PInches(4.0), PInches(2.6),
          "Home",
          "Dashboard with accent-striped feature cards, hero banner "
          "and one-click quick actions.")
    _card(s, PInches(4.7), PInches(1.3), PInches(4.0), PInches(2.6),
          "Translator",
          "Live ASCII ↔ Morse, history panel, .txt + .wav export, "
          "tone-frequency and unit sliders.",
          accent=ACCENT_GREEN)
    _card(s, PInches(8.8), PInches(1.3), PInches(3.9), PInches(2.6),
          "Cheat Sheet",
          "Click any of the 40+ symbol cards to hear the Morse — "
          "perfect for memorising during viva prep.",
          accent=ACCENT_RED)
    _card(s, PInches(0.6), PInches(4.0), PInches(4.0), PInches(2.6),
          "ASCII + IVT",
          "Full ASCII table with Morse codes plus the IVT entries "
          "our DOS calls use — printable for study notes.")
    _card(s, PInches(4.7), PInches(4.0), PInches(4.0), PInches(2.6),
          "COAL Quiz",
          "Twelve interactive questions with instant green/red "
          "feedback and a final score summary.",
          accent=ACCENT_GREEN)
    _card(s, PInches(8.8), PInches(4.0), PInches(3.9), PInches(2.6),
          "Polish",
          "Starts maximised, tooltips on every control, keyboard "
          "shortcuts, welcome overlay for new users.",
          accent=ACCENT_RED)
    _footer(s, 15, total_slides)

    # ----- Slide 16: Challenges & Ethics ------------------------------
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _title_strip(s, "Opportunities, Challenges & Ethics")
    _card(s, PInches(0.7), PInches(1.3), PInches(6.0), PInches(2.7),
          "Opportunities",
          "Battery-friendly emergency beacons, accessibility "
          "tool for speech-impaired users, amateur-radio training.",
          accent=ACCENT_GREEN)
    _card(s, PInches(6.9), PInches(1.3), PInches(5.7), PInches(2.7),
          "Challenges",
          "Portability across MASM/TASM/emu8086 dialects; producing "
          "clean tones without direct port access.",
          accent=ACCENT_RED)
    _card(s, PInches(0.7), PInches(4.1), PInches(11.9), PInches(2.7),
          "Ethics",
          "Respect frequency allocations, never transmit false distress "
          "calls, always self-identify on air and use the technology to "
          "help — never to deceive or eavesdrop.")
    _footer(s, 16, total_slides)

    # ----- Slide 17: Conclusion ---------------------------------------
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _title_strip(s, "Conclusion & Future Work")
    _add_bullets(s, [
        "Three-layer implementation (assembly + simulator + GUI) "
        "cleanly maps the entire CS-234 syllabus.",
        "Assembly program demonstrates every major addressing mode.",
        "CPU Simulator makes every instruction observable on screen.",
        "Desktop GUI turns the program into a presentable demo with "
        "audio, animation, cheat sheet, ASCII/IVT reference and quiz.",
        "Full documentation — proposal, IEEE report, viva guide.",
        "Future: microcontroller beacon with LED, audio decoder, "
        "radio integration, multilingual character support.",
    ], PInches(0.7), PInches(1.5), PInches(12), PInches(5), size=17)
    _footer(s, 17, total_slides)

    # ----- Slide 18: Thanks -------------------------------------------
    s = prs.slides.add_slide(blank); _slide_bg(s)
    _add_rect(s, PInches(0), PInches(3.1), PInches(13.333),
              PInches(0.06), CYAN)
    _add_text(s, "Thank you!",
              PInches(0.8), PInches(2.0), PInches(12), PInches(1.2),
              size=72, bold=True, color=TEXT_LIGHT,
              align=PP_ALIGN.CENTER)
    _add_text(s, "Questions  ·  Feedback  ·  Discussion",
              PInches(0.8), PInches(3.35), PInches(12), PInches(0.6),
              size=22, color=CYAN, align=PP_ALIGN.CENTER)
    _add_text(s, "Mustafa Shahid   |   BSCS-14B   |   CMS 500889",
              PInches(0.8), PInches(4.5), PInches(12), PInches(0.5),
              size=18, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    _add_text(s, "CS-234  Computer Organization & Assembly Language",
              PInches(0.8), PInches(5.0), PInches(12), PInches(0.5),
              size=15, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    out = DOC_DIR / "03_Presentation.pptx"
    prs.save(out)
    return out


# ==========================================================================
# 4. VIVA PREPARATION GUIDE
# ==========================================================================
def build_viva_guide() -> Path:
    doc = Document()
    set_default_font(doc, size=11)
    for sec in doc.sections:
        sec.top_margin = Inches(0.8)
        sec.bottom_margin = Inches(0.8)
        sec.left_margin = Inches(0.9)
        sec.right_margin = Inches(0.9)

    # Title block
    t = doc.add_paragraph()
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    styled_run(t, "VIVA PREPARATION GUIDE", size=22, bold=True,
               color="0F172A")
    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    styled_run(sub, "Morse Code Communication & Emergency Signaling System",
               size=12, bold=True, color="0EA5E9")
    sub2 = doc.add_paragraph()
    sub2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    styled_run(sub2,
               "Mustafa Shahid  ·  BSCS-14B  ·  CMS 500889  ·  CS-234 COAL",
               size=11, color="475569")
    add_horizontal_line(sub2)

    add_body(
        doc,
        "This concise guide prepares you for the viva.  Memorise the "
        "30-second pitch, skim every “Likely Question → Short Answer” "
        "pair, and you will comfortably handle anything the faculty asks.",
        italic=True, size=10,
    )

    # -----------------------------------------------------------------
    add_heading(doc, "1.  30-Second Elevator Pitch", level=2,
                color="0EA5E9")
    add_body(
        doc,
        "“My project is a Morse code communication and emergency "
        "signaling system.  It has two parts.  The first is an 8086 "
        "assembly program that converts text into Morse using a lookup "
        "table accessed through indexed addressing, and uses DOS "
        "interrupts for input/output.  The second is a modern Windows "
        "desktop application that mirrors the same logic, plays audible "
        "beeps, visualises the dots and dashes live and includes an "
        "emergency SOS broadcaster.  Together they cover the core COAL "
        "topics — addressing modes, interrupts, the IVT, "
        "arithmetic/logical operations and I/O — in a real-world "
        "application used in aviation, maritime rescue and amateur "
        "radio.”",
        italic=False,
    )

    # -----------------------------------------------------------------
    add_heading(doc, "2.  Key Facts To Remember", level=2, color="0EA5E9")
    add_bullets(doc, [
        "Morse code: dots and dashes; standard is ITU-R M.1677-1.",
        "Timing units  →  dot 1, dash 3, intra-char gap 1, "
        "inter-char gap 3, inter-word gap 7.",
        "Lookup table  →  each entry is 6 bytes (5 symbols + 1 length).",
        "Address formula:  base + (index × 6).",
        "Letter index  =  ASCII – 'A';  digit index = ASCII – '0' + 26.",
        "Case conversion trick:  AND AL, 0DFh  (clears bit 5).",
        "DOS services used:  09h print string, 0Ah buffered input, "
        "02h print char, 4Ch terminate.",
        "SOS  =  ... --- ...   (3 dots, 3 dashes, 3 dots).",
        "BEL character 07h triggers the PC speaker.",
        "GUI uses Python’s winsound.Beep(frequency, duration).",
    ])

    # -----------------------------------------------------------------
    add_heading(doc, "3.  Likely Questions and Short Answers",
                level=2, color="0EA5E9")

    qa = [
        ("Why did you choose Morse code?",
         "It is simple to explain, has enduring real-world impact "
         "(aviation, maritime rescue, amateur radio, accessibility), "
         "and its lookup-table nature exercises indexed and "
         "base-plus-offset addressing modes elegantly."),

        ("What COAL concepts does your project cover?",
         "Data definition directives (DB, EQU), direct, indirect, "
         "indexed and base+offset addressing, arithmetic and logical "
         "operations, software interrupts, the IVT, procedures, and "
         "low-level I/O through the PC speaker."),

        ("How is your Morse lookup table structured?",
         "It is a flat byte-array of fixed 6-byte records.  Bytes 0–4 "
         "hold the Morse symbols, space-padded to width 5; byte 5 "
         "stores the code length."),

        ("How do you compute the address of a letter's entry?",
         "addr = morse_table + (index × 6).  Index for 'A'–'Z' is "
         "ASCII – 'A'; for '0'–'9' it is ASCII – '0' + 26."),

        ("Explain indexed addressing with a line from your code.",
         "MOV DL, [BX + DI].  BX is the table base, DI is the record "
         "offset (index × 6), so the CPU computes the effective "
         "address at run time."),

        ("What is base-plus-offset addressing?  Example?",
         "MOV CH, [BX + DI + 5].  Effective address = BX + DI + 5. "
         "We use it to reach the length byte inside a record."),

        ("How do you convert lowercase to uppercase?",
         "With a single logical operation:  AND AL, 0DFh.  This "
         "clears bit 5 of the ASCII code, which is the difference "
         "between a lowercase and an uppercase letter."),

        ("What is the Interrupt Vector Table (IVT)?",
         "A table of 256 four-byte entries starting at physical "
         "address 0.  Each INT n instruction looks up entry n to "
         "find the segment:offset of the handler and jumps there."),

        ("Which DOS interrupts do you use, and why?",
         "INT 21h with functions 09h (print string), 0Ah (buffered "
         "input), 02h (print character) and 4Ch (terminate).  They "
         "isolate the program from hardware specifics."),

        ("How do you make the PC speaker beep?",
         "In the assembly version we print the BEL character (07h) "
         "through INT 21h/02h.  In the desktop app we call "
         "winsound.Beep(frequency, duration) for cleaner tones."),

        ("What happens when the user types an unsupported character?",
         "It is silently skipped by the CMP ... / JB NEXT_CHAR "
         "guards, so malformed input never halts the demo."),

        ("How is an SOS broadcast implemented?",
         "After translation, MAIN calls PLAY_SOS which emits three "
         "short beeps, three long beeps, three short beeps — the "
         "internationally recognised distress signal."),

        ("How do short and long beeps differ?",
         "By duration: a dash is three times longer than a dot, per "
         "ITU-R M.1677-1.  Our delay routine multiplies the base "
         "wait loop by three for dashes."),

        ("What is the Morse code of your CMS ID, 500889?",
         "..... ----- ----- ---.. ---.. ----."),

        ("What is the Morse code for SOS?",
         "... --- ..."),

        ("Why is Morse still used today?",
         "Because it survives extreme noise, requires only on/off "
         "signalling, works with milliwatts of power and is "
         "decodable by the human ear alone."),

        ("What is the difference between direct and indirect "
         "addressing?",
         "Direct addressing names an absolute memory location in the "
         "instruction, e.g. MOV AL, [input_buffer].  Indirect "
         "addressing uses a register as a pointer, e.g. MOV AL, [SI]."),

        ("What is CALL / RET doing under the hood?",
         "CALL pushes the return address onto the stack and jumps to "
         "the procedure.  RET pops the saved address and returns, "
         "so the stack provides re-entrancy for nested calls."),

        ("Why did you write the GUI in Python?",
         "Python + CustomTkinter gives a beautiful cross-platform UI "
         "quickly, ships winsound for clean audio and keeps the demo "
         "focused on the COAL concepts rather than UI plumbing."),

        ("Can you run the assembly program?",
         "Yes — open Assembly/MorseTranslator.asm in emu8086, press "
         "\"Emulate\" and then \"Run\".  The banner appears, it "
         "prompts for input, prints the Morse and beeps SOS."),

        ("What is the purpose of the PUSH/POP pairs around "
         "BEEP_SHORT?",
         "They save AX and DX before the BIOS call and restore them "
         "afterwards, preserving the caller’s context — this is "
         "textbook stack discipline."),

        ("Why did you use RECORD_SIZE EQU 6?",
         "EQU defines a named compile-time constant.  Using a "
         "symbolic name makes the code self-documenting and avoids "
         "magic numbers."),

        ("How would you extend this project for real hardware?",
         "Port the translator to an 8-bit microcontroller, connect "
         "an LED + buzzer for output, and power it from a coin "
         "cell — giving a battery-powered SOS beacon."),

        ("What are the ethical considerations?",
         "Never broadcast false distress calls, respect licensed "
         "frequency allocations, always identify yourself on air "
         "and use the technology to help people, not to deceive."),

        ("Walk me through your CPU Simulator.",
         "It is a page inside the desktop app.  I type any letter or "
         "digit and click Load — the simulator builds the exact "
         "instruction trace the assembly code issues for that "
         "character.  Step executes one instruction at a time, the "
         "ten registers update live in hex, and the memory panel "
         "highlights the lookup-table record being accessed through "
         "indexed addressing."),

        ("How does your application export Morse as a .wav file?",
         "Pure Python: I compute the number of samples for each dot "
         "and dash using the current frequency and unit duration, "
         "generate sine-wave samples with a short attack/decay fade "
         "to avoid clicks, pack them as 16-bit little-endian PCM and "
         "write them with the wave module. No third-party audio "
         "library is needed."),

        ("Why fourteen pages instead of a single big window?",
         "Each page has a clear, single purpose — translator, "
         "telegraph key, trainer, simulator, cheat sheet, quiz, "
         "settings and so on — which reduces cognitive load during "
         "the demo and gives the examiner a logical path to ask "
         "questions about any specific COAL topic."),

        ("Explain your Telegraph Key page.",
         "It measures how long I press the button (or the Space key). "
         "A press shorter than twice the current unit duration is a "
         "dot, anything longer is a dash.  When I pause for three "
         "units the app commits the accumulated dots and dashes as a "
         "letter; a longer pause inserts a word space.  It decodes "
         "my rhythm into plain text live — exactly like a human "
         "operator on a telegraph line."),

        ("What does your Morse Trainer do?",
         "It is a gamified decode challenge.  I pick a difficulty "
         "(Letter, Word or Phrase), the app shows a random Morse "
         "pattern, I type the plain text and press Enter.  Correct "
         "answers earn points based on speed; three correct answers "
         "in a row multiply my score.  It reinforces the lookup "
         "direction that my assembly program performs."),

        ("What’s new in the CPU Simulator?",
         "Beyond single-step and Run-All it now has a Play / Pause "
         "button with a speed slider between 150 ms and 1500 ms per "
         "instruction, and a FLAGS panel that breaks the status "
         "register into its six bits — Zero, Carry, Sign, Overflow, "
         "Parity and Auxiliary — each illuminated when it is set."),

        ("What’s the difference between your GUI and the assembly "
         "program?",
         "Behavioural parity: both convert text to Morse using the "
         "same 6-byte lookup table and the same index formula.  The "
         "assembly program uses the PC speaker via BEL; the GUI uses "
         "winsound.Beep for cleaner tones and adds visualisation, "
         "history, audio export and educational panels."),
    ]

    for q, a in qa:
        pq = doc.add_paragraph()
        pq.paragraph_format.space_before = Pt(4)
        pq.paragraph_format.space_after = Pt(1)
        styled_run(pq, "Q.  ", size=11, bold=True, color="0EA5E9")
        styled_run(pq, q, size=11, bold=True, color="0F172A")
        pa = doc.add_paragraph()
        pa.paragraph_format.space_after = Pt(6)
        pa.paragraph_format.left_indent = Inches(0.28)
        styled_run(pa, "A.  ", size=11, bold=True, color="22C55E")
        styled_run(pa, a, size=11, color="334155")

    # -----------------------------------------------------------------
    add_heading(doc, "4.  Demo Flow During Viva", level=2, color="0EA5E9")
    add_bullets(doc, [
        "Double-click Run_Application.bat — the window opens maximised.",
        "Home page → briefly introduce the feature cards.",
        "Translator → type HELLO WORLD, click ▶ Play Audio.",
        "Save .txt and .wav → mention how WAV is generated from "
        "sine-wave samples in pure Python.",
        "Cheat Sheet → click any letter card to hear it.",
        "CPU Simulator → type 'S', click Load, click Step repeatedly; "
        "narrate each register update and each indexed fetch.",
        "Live Visualiser → transmit SOS HELP.",
        "Emergency SOS → BROADCAST SOS.",
        "ASCII/IVT Reference and COAL Concepts → walk the cards.",
        "COAL Quiz → answer one question to demonstrate feedback.",
        "Assembly Code page → narrate the key sections of the .asm.",
        "If asked, also run the .asm in emu8086 (File → Open → Emulate).",
    ])

    # -----------------------------------------------------------------
    add_heading(doc, "5.  Recovery Lines If You Get Stuck",
                level=2, color="0EA5E9")
    add_bullets(doc, [
        "“Let me show you directly in the code — can I open the "
        "Assembly Code page?”",
        "“That’s a great question.  The key idea is that the Morse "
        "table has fixed-size records, which lets us use indexed "
        "addressing to look up any character in constant time.”",
        "“If you don’t mind, let me run the demo first and then I’ll "
        "walk through that part.”",
    ])

    out = DOC_DIR / "04_Viva_Preparation_Guide.docx"
    doc.save(out)
    return out


# ==========================================================================
# MAIN
# ==========================================================================
def main() -> None:
    print("Building Project Proposal ...")
    p1 = build_proposal()
    print(f"   ->  {p1}")

    print("Building Project Report (IEEE) ...")
    p2 = build_report()
    print(f"   ->  {p2}")

    print("Building PowerPoint presentation ...")
    p3 = build_pptx()
    print(f"   ->  {p3}")

    print("Building Viva preparation guide ...")
    p4 = build_viva_guide()
    print(f"   ->  {p4}")

    print("\nAll documents generated successfully.")


if __name__ == "__main__":
    main()
